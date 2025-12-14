"""PPTX file processor."""

from pptx import Presentation


def process_pptx(file_path: str) -> str:
    """
    Extract text from PPTX file.
    
    Args:
        file_path: Path to PPTX file
        
    Returns:
        Extracted text content
    """
    prs = Presentation(file_path)
    text = ""
    
    for slide_num, slide in enumerate(prs.slides, 1):
        text += f"\n--- Slide {slide_num} ---\n"
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            
            # Extract text from tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
    
    return text
